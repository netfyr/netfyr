#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# White background theme with black + red text
BG = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
RED = RGBColor(0xC0, 0x1B, 0x2B)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
MID_GRAY = RGBColor(0x77, 0x77, 0x77)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)
BOX_BG = RGBColor(0xFA, 0xFA, 0xFA)
BOX_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
RED_BOX_BG = RGBColor(0xFD, 0xF0, 0xF0)
RED_BOX_BORDER = RGBColor(0xC0, 0x1B, 0x2B)
DARK_GREEN = RGBColor(0x1B, 0x7A, 0x3D)
GREEN_BG = RGBColor(0xF0, 0xFA, 0xF4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_box(slide, left, top, width, height, items, font_size=16,
                   color=DARK_GRAY, bullet_color=RED, spacing=8):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(spacing)
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=RED, body_color=DARK_GRAY, bg=BOX_BG, border=BOX_BORDER,
             body_font_size=13):
    add_rounded_rect(slide, left, top, width, height, bg, border)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1),
                width - Inches(0.4), Inches(0.4), title,
                font_size=15, color=title_color, bold=True)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.5),
        width - Inches(0.4), height - Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(body_font_size)
        p.font.color.rgb = body_color
        p.font.name = "Calibri"
        p.space_after = Pt(2)


# =====================================================================
# SLIDE 1: Title — Problem + Approach
# =====================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, BG)

add_textbox(slide, Inches(1), Inches(0.7), Inches(11), Inches(1.2),
            "netfyr", font_size=60, color=RED, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_textbox(slide, Inches(1), Inches(1.9), Inches(11), Inches(0.8),
            "Declarative Network Configuration for Linux",
            font_size=28, color=BLACK, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(4), Inches(2.8), Inches(5.3), Pt(2), RED)

# Problem column
add_rounded_rect(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(3.9), RED_BOX_BG, RED_BOX_BORDER)
add_textbox(slide, Inches(1), Inches(3.3), Inches(5), Inches(0.4),
            "The Problem", font_size=20, color=RED, bold=True)
problem_items = [
    "Existing solutions rely on always-running daemons, even for static or simple configurations",
    "No visibility into why network state changed",
    "No visibility on previous network states",
    "Multiple configuration sources can target the same interfaces, causing silent conflicts",
]
add_bullet_box(slide, Inches(1), Inches(3.9), Inches(5.1), Inches(3),
               problem_items, font_size=17, color=DARK_GRAY, spacing=10)

# Solution column
add_rounded_rect(slide, Inches(7), Inches(3.2), Inches(5.5), Inches(3.9), GREEN_BG, DARK_GREEN)
add_textbox(slide, Inches(7.2), Inches(3.3), Inches(5), Inches(0.4),
            "netfyr's Approach", font_size=20, color=DARK_GREEN, bold=True)
solution_items = [
    "Flexible architecture: daemon or daemonless operation with reduced functionality",
    "Track the provenance of every network configuration change",
    "Append-only journal recording all state transitions",
    "Reconciliation of multiple configurations with per-field priority and conflict detection",
]
add_bullet_box(slide, Inches(7.2), Inches(3.9), Inches(5.1), Inches(3),
               solution_items, font_size=17, color=DARK_GRAY, bullet_color=DARK_GREEN, spacing=10)


# =====================================================================
# SLIDE 2: Architecture — Component blocks with dependencies
# =====================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Architecture", font_size=36, color=RED, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.5),
            "Modular components with layered dependencies",
            font_size=16, color=MID_GRAY)

# Layout: 4 rows of component blocks
# Row 0 (top): CLI, Daemon
# Row 1: Backend, Journal, Varlink
# Row 2: Policy, Reconcile
# Row 3 (bottom): State

block_h = Inches(1.15)
row_gap = Inches(0.35)
col_w = Inches(3.6)
col_gap = Inches(0.3)

# Compute row tops
row0_top = Inches(1.5)
row1_top = row0_top + block_h + row_gap
row2_top = row1_top + block_h + row_gap
row3_top = row2_top + block_h + row_gap

# Center positions for different row widths
full_width = Inches(12)
margin_left = Inches(0.667)

# Row 0: 2 blocks centered
r0_total = 2 * col_w + col_gap
r0_left = margin_left + (full_width - r0_total) / 2
blocks_row0 = [
    ("CLI", "7 subcommands: apply, query,\nhistory, revert, show, diagnose, completions", r0_left),
    ("Daemon", "Long-running service with Varlink server,\nDHCP factory management, netlink monitor", r0_left + col_w + col_gap),
]

# Row 1: 3 blocks centered
r1_total = 3 * col_w + 2 * col_gap
r1_left = margin_left + (full_width - r1_total) / 2
blocks_row1 = [
    ("Backend", "NetworkBackend trait with rtnetlink\nimplementation; query and apply operations", r1_left),
    ("Journal", "Append-only NDJSON history with\nprovenance tracking and state snapshots", r1_left + col_w + col_gap),
    ("Varlink", "IPC protocol for CLI-daemon\ncommunication over Unix sockets", r1_left + 2 * (col_w + col_gap)),
]

# Row 2: 2 blocks centered
r2_total = 2 * col_w + col_gap
r2_left = margin_left + (full_width - r2_total) / 2
blocks_row2 = [
    ("Policy", "Policy model with static and dynamic\n(DHCPv4) factories; YAML loading", r2_left),
    ("Reconcile", "Per-field priority merge with explicit\nconflict detection and diff generation", r2_left + col_w + col_gap),
]

# Row 3: 1 block centered
r3_total = col_w
r3_left = margin_left + (full_width - r3_total) / 2
blocks_row3 = [
    ("State", "Core types: State, Selector, Value;\nJSON Schema validation; StateSet operations", r3_left),
]

all_rows = [
    (row0_top, blocks_row0),
    (row1_top, blocks_row1),
    (row2_top, blocks_row2),
    (row3_top, blocks_row3),
]

for row_top, blocks in all_rows:
    for title, desc, left in blocks:
        add_rounded_rect(slide, left, row_top, col_w, block_h, BOX_BG, RED)
        add_textbox(slide, left + Inches(0.2), row_top + Inches(0.08),
                    col_w - Inches(0.4), Inches(0.35), title,
                    font_size=16, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
        # Description text
        txBox = slide.shapes.add_textbox(left + Inches(0.15), row_top + Inches(0.45),
                                         col_w - Inches(0.3), Inches(0.65))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(desc.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)

# Draw arrows between rows
arrow_color = RGBColor(0xBB, 0xBB, 0xBB)
arrow_rows = [
    (row0_top + block_h, row1_top),
    (row1_top + block_h, row2_top),
    (row2_top + block_h, row3_top),
]
for arrow_top, arrow_bottom in arrow_rows:
    mid = margin_left + full_width / 2
    h = arrow_bottom - arrow_top
    add_textbox(slide, mid - Inches(0.25), arrow_top, Inches(0.5), h,
                "▼", font_size=14, color=arrow_color, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 3: Key Features (concise)
# =====================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Key Features", font_size=36, color=RED, bold=True)

# 3x2 grid of cards
cols = 3
rows = 2
card_w = Inches(3.7)
card_h = Inches(2.5)
start_left = Inches(0.8)
start_top = Inches(1.2)
h_gap = Inches(0.2)
v_gap = Inches(0.2)

features = [
    ("DHCPv4 Dynamic Factory",
     ["Dual-socket: AF_PACKET before lease,",
      "  AF_INET (UDP) after lease",
      "Full lifecycle: Discover → Offer →",
      "  Request → Ack → Renew → Rebind",
      "Leases participate in priority merge",
      "  with static policies"]),
    ("History Journal",
     ["Append-only NDJSON format",
      "Triggers: PolicyApply, DhcpEvent,",
      "  ExternalChange, DaemonStartup, Revert",
      "Rotation at 10k entries or 50MB (gzip)",
      "Configurable retention (default 90 days)",
      "Revert to any sequence ID"]),
    ("External Change Detection",
     ["Netlink monitoring (link, addr, route)",
      "500ms debounce window",
      "Does NOT auto-revert:",
      "  changes journaled for operator review"]),
    ("Varlink API (daemon)",
     ["Read-only (any user): Query, DryRun,",
      "  GetStatus, GetHistory, GetJournalEntry,",
      "  GetShowInfo",
      "Write (root only): SubmitPolicies, Revert"]),
    ("Device Selection",
     ["Multi-field AND matching:",
      "  name, entity_type, driver,",
      "  mac, pci_path, labels"]),
    ("Schema Validation",
     ["Inheritance via x-netfyr-inherit:",
      "  ethernet ← ip.json + link.json",
      "Custom attributes:",
      "  x-netfyr-writable,",
      "  x-netfyr-keep-when-absent,",
      "  x-netfyr-comparison-keys"]),
]

for idx, (title, lines) in enumerate(features):
    col = idx % cols
    row = idx // cols
    left = start_left + col * (card_w + h_gap)
    top = start_top + row * (card_h + v_gap)
    add_card(slide, left, top, card_w, card_h, title, lines)


# =====================================================================
# SLIDE 4: Tech Stack, Data Model, Testing
# =====================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Tech Stack & Data Model", font_size=36, color=RED, bold=True)

# Left column: tech stack
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(4.3), Inches(5.8), BOX_BG, BOX_BORDER)
add_textbox(slide, Inches(1), Inches(1.3), Inches(3), Inches(0.35),
            "Technology Stack", font_size=17, color=RED, bold=True)

stack = [
    ("Language", "Rust (Edition 2021)"),
    ("Async", "Tokio 1.x"),
    ("Netlink", "rtnetlink 0.20"),
    ("DHCP", "dhcproto 0.14"),
    ("IPC", "Varlink / Unix socket"),
    ("CLI", "clap 4 (derive)"),
    ("Schema", "jsonschema (v2020-12)"),
    ("Journal", "NDJSON + flate2"),
    ("Logging", "tracing"),
    ("Systemd", "sd-notify, Type=notify"),
]

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.75), Inches(4), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (key, val) in enumerate(stack):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p.add_run()
    r1.text = f"{key:<12}"
    r1.font.size = Pt(14)
    r1.font.color.rgb = RED
    r1.font.name = "Consolas"
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(14)
    r2.font.color.rgb = BLACK
    r2.font.name = "Calibri"
    p.space_after = Pt(5)

# Center column: core data model
add_rounded_rect(slide, Inches(5.3), Inches(1.2), Inches(3.7), Inches(5.8), BOX_BG, BOX_BORDER)
add_textbox(slide, Inches(5.5), Inches(1.3), Inches(3.3), Inches(0.35),
            "Core Data Model", font_size=17, color=RED, bold=True)

model_lines = [
    ("State", "entity_type, selector, fields,"),
    ("", "  metadata (UUIDv7), priority, policy_ref"),
    ("Value", "String | U64 | I64 | Bool |"),
    ("", "  IpAddr | IpNetwork | List | Map"),
    ("Provenance", "UserConfigured | KernelDefault |"),
    ("", "  ExternalTool | Derived"),
    ("Policy", "name, factory_type (Static | Dhcpv4),"),
    ("", "  priority (default 100), selector, states"),
    ("Selector", "name?, type?, driver?, mac?,"),
    ("", "  pci_path?, labels? — AND matching"),
]

txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.75), Inches(3.3), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (label, desc) in enumerate(model_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    if label:
        r1 = p.add_run()
        r1.text = label + "  "
        r1.font.size = Pt(12)
        r1.font.color.rgb = RED
        r1.font.name = "Consolas"
        r1.font.bold = True
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(12)
    r2.font.color.rgb = DARK_GRAY
    r2.font.name = "Calibri"
    p.space_after = Pt(2)

# Right column: testing & CLI
right = Inches(9.2)
add_card(slide, right, Inches(1.2), Inches(3.8), Inches(2.5),
         "Testing",
         ["130+ integration test scripts (shell)",
          "Unprivileged: unshare --user --net",
          "veth pairs + dnsmasq for DHCP tests",
          "No-skip policy: fail prerequisites → exit(1)",
          "make integration-test SPEC=NNN",
          "End-to-end: apply, query, history, revert,",
          "  DHCP, conflicts, external changes"])

add_card(slide, right, Inches(3.9), Inches(3.8), Inches(1.7),
         "CLI Subcommands",
         ["apply, query, history, revert,",
          "show, diagnose, completions",
          "--dry-run preview on apply & revert",
          "Exit: 0 success, 1 partial, 2 fatal"])

add_card(slide, right, Inches(5.8), Inches(3.8), Inches(1.2),
         "Packaging",
         ["RPM spec with systemd units",
          "Man pages via clap_mangen (xtask)",
          "License: Apache 2.0"])


prs.save("/home/test/work/netfyr/netfyr_presentation.pptx")
print("Saved netfyr_presentation.pptx")
